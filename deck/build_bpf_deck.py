#!/usr/bin/env python3
"""Fill the Builders Pitch Fest 2026 template with Navixa's submission content.

Keeps the template's own layout, branding and question prompts; writes the
answers underneath each question. Run:  python3 deck/build_bpf_deck.py
"""
import copy
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = sys.argv[1] if len(sys.argv) > 1 else "template.pptx"
OUTPUT = sys.argv[2] if len(sys.argv) > 2 else "BPF2026_PrototypeSubmission_PitchDeck_Navixa.pptx"

INK = RGBColor(0x21, 0x21, 0x21)
BODY = RGBColor(0x3C, 0x3C, 0x3C)
MUTED = RGBColor(0x5E, 0x5E, 0x5E)
ORANGE = RGBColor(0xE1, 0x66, 0x3B)

ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def _pPr(p):
    return p._p.get_or_add_pPr()


def no_bullet(p):
    pPr = _pPr(p)
    for tag in ('buChar', 'buAutoNum', 'buNone'):
        for e in pPr.findall(ns + tag):
            pPr.remove(e)
    pPr.append(pPr.makeelement(ns + 'buNone', {}))


def dot_bullet(p):
    pPr = _pPr(p)
    for tag in ('buChar', 'buAutoNum', 'buNone'):
        for e in pPr.findall(ns + tag):
            pPr.remove(e)
    f = pPr.makeelement(ns + 'buFont', {'typeface': 'Arial'})
    c = pPr.makeelement(ns + 'buChar', {'char': '•'})
    pPr.append(f)
    pPr.append(c)


def line(tf, text, *, size=10, bold=False, colour=BODY, indent=0, space_before=0,
         first=False, bullet=False, italic=False):
    """Append one paragraph. `indent` is in inches from the box's left edge."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    pPr = _pPr(p)
    pPr.set('marL', str(Emu(Pt(0) + int(indent * 914400))))
    pPr.set('indent', str(-114300 if bullet else 0))
    if bullet:
        dot_bullet(p)
    else:
        no_bullet(p)
    p.space_before = Pt(space_before)
    p.space_after = Pt(0)
    p.line_spacing = 1.12
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = 'Arial'
    r.font.color.rgb = colour
    return p


def body_of(slide):
    """The template's answer box is the widest text shape on the slide."""
    cands = [s for s in slide.shapes if s.has_text_frame and s.width > Emu(5000000)]
    return max(cands, key=lambda s: s.height)


def q(tf, n, text, first=False):
    line(tf, f'{n}. {text}', size=10.5, bold=True, colour=INK,
         space_before=0 if first else 11, first=first)


def a(tf, text, indent=0.28, size=10, colour=BODY, space_before=3, bold=False):
    return line(tf, text, size=size, colour=colour, indent=indent,
                space_before=space_before, bold=bold)


def b(tf, text, indent=0.42, size=10, space_before=2):
    line(tf, text, size=size, colour=BODY, indent=indent,
         space_before=space_before, bullet=True)


prs = Presentation(TEMPLATE)
S = prs.slides

# ---------------------------------------------------------------- slide 3
tf = body_of(S[2]).text_frame
tf.clear()
q(tf, 1, 'What does your startup do', first=True)
a(tf, 'Navixa is an AI career-navigation web app for students and early-career professionals. '
      'It replaces the six-to-eight disconnected tools a job hunt normally needs with one workspace '
      'where every feature shares the same context: live job search, AI resume tailoring against a '
      'specific posting, an application tracker with funnel analytics, and job-specific interview '
      'practice with spoken answers. Free to use, runs entirely in the browser.')
q(tf, 2, 'What milestone best represents your progress so far?')
a(tf, 'A fully deployed, working product — not a prototype or mockup. Live at navixa-woad.vercel.app '
      'with Google sign-in, cloud sync and an admin console.')
b(tf, 'Five feature phases shipped end-to-end: resume tailoring, tracker analytics, saved searches, '
      'interview practice, shareable public profiles.')
b(tf, '114 automated tests run green (33 functional, 62 security, 19 rate-limiting), plus a documented '
      'security audit that found and fixed two exploitable vulnerabilities.')
b(tf, 'Validated informally with a small group of peer testers; no paid users yet.')

# ---------------------------------------------------------------- slide 4
tf = body_of(S[3]).text_frame
tf.clear()
q(tf, 1, 'What problem are you solving, and who experiences it most acutely?', first=True)
a(tf, 'A job search is run across tools that cannot see each other: job boards for listings, Word or Canva '
      'for the resume, a spreadsheet for tracking, YouTube for interview prep, and a chatbot for advice. '
      'The person is the integration layer. The resume never knows the posting, the AI never knows the resume, '
      'and nothing remembers what already worked.')
a(tf, 'Felt most acutely by first-time job seekers — final-year students and those 0–3 years in — who have no '
      'placement support, no mentor, and no way to tell whether their application is weak or the market is.')
q(tf, 2, 'What evidence validates that this is a meaningful problem worth solving?')
b(tf, 'India Skills Report 2026: only 56.35% of graduates are assessed as employable — close to half enter the '
      'market underprepared.')
b(tf, 'AISHE 2023–24: 45 million students are enrolled in Indian higher education, 10.2 million in STEM.')
b(tf, 'Peer testers independently described the same fragmented workflow, and the same complaint about generic '
      'AI advice that ignores their actual resume.')

# ---------------------------------------------------------------- slide 5
tf = body_of(S[4]).text_frame
tf.clear()
q(tf, 1, 'Who is your ideal customer, and who makes the buying decision?', first=True)
a(tf, 'Primary user: final-year students and professionals 0–3 years into their career in India, running an '
      'active search without institutional support. In this B2C motion the user is also the buyer, so the price '
      'point must stay low.')
a(tf, 'Second motion (B2B2C): college placement cells and training institutes, where the placement officer or '
      'training head buys and students use. They already carry placement-rate targets and today track outcomes '
      'in spreadsheets.')
q(tf, 2, 'How large is the opportunity you are targeting?')
b(tf, 'TAM — 45M students in Indian higher education (AISHE 2023–24).')
b(tf, 'SAM — ~10.2M STEM enrolments, the segment with the most structured, competitive hiring.')
b(tf, 'Beachhead — placement cells of engineering colleges; 48,246 colleges operate nationally, and a single '
      'institutional sale reaches hundreds of students at once.')

# ---------------------------------------------------------------- slide 6
tf = body_of(S[5]).text_frame
tf.clear()
q(tf, 1, 'What is your solution, and how does it solve the identified problem?', first=True)
a(tf, 'One workspace where the features share context instead of sitting side by side. Tailoring reads the '
      'tracker; interview prep reads both the resume and the job posting. The integration work the user was '
      'doing by hand is what the product does.')
q(tf, 2, 'What are the core capabilities of your product?')
b(tf, 'Job search across four public boards, de-duplicated, with skill-match scoring and salary insight.')
b(tf, 'AI resume tailoring per posting — match score, missing keywords, bullet rewrites — plus cover letters, '
      'skill-gap analysis and PDF/Word resume import.')
b(tf, 'Application tracker with funnel analytics, stale-application nudges, AI follow-up drafts, CSV and '
      'calendar export.')
b(tf, 'Interview prep that generates job-specific questions, including deliberate "gap probe" questions aimed '
      'at weaknesses visible between the resume and the posting, with spoken practice and delivery metrics.')
q(tf, 3, 'What measurable value does your solution deliver?')
b(tf, 'Tailoring a resume to a posting drops from roughly 30 minutes of manual work to under a minute.')
b(tf, 'Interview questions target the candidate’s actual gaps rather than a generic list.')
b(tf, 'Every AI feature has a deterministic fallback, so the product still works when the model is unavailable.')

# ---------------------------------------------------------------- slide 7  (+ diagram)
tf = body_of(S[6]).text_frame
tf.clear()
q(tf, 1, 'How is your solution architecture end-to-end, including data flow and major system components?',
  first=True)
a(tf, 'Browser SPA (vanilla ES modules, no build step) → three stateless serverless functions on Vercel '
      '→ Supabase Postgres. Job data is pulled from public boards through an allowlisted same-origin proxy. '
      'Every user’s data is isolated by Postgres row-level security: app state lives in a private row that '
      'only its owner can read — not even an admin can.')
q(tf, 2, 'What role does AI play within your solution, and which capabilities or workflows are powered by it?')
a(tf, 'AI powers resume parsing from PDF/DOCX, per-posting tailoring and bullet rewrites, cover letters, '
      'skill-gap analysis, interview-question generation and answer feedback. It is used where judgement is '
      'needed — not for everything: delivery metrics (words per minute, filler-word rate) are computed '
      'locally so they work offline and cost nothing, and every AI path falls back to a deterministic result if '
      'the model is unreachable.')

sl = S[6]
box_w, gap = Emu(1560000), Emu(150000)
x0, y = Emu(430000), Emu(3960000)
steps = [('Browser SPA', 'resume · tracker\ninterview prep'),
         ('Serverless API', 'LLM relay · proxy\nrate-limited'),
         ('LLM gateway', 'model-agnostic\nwith fallbacks'),
         ('Job boards', '4 public APIs\nde-duplicated'),
         ('Supabase', 'Postgres + RLS\nper-user isolation')]
for i, (head, sub) in enumerate(steps):
    x = x0 + i * (box_w + gap)
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, Emu(620000))
    sh.adjustments[0] = 0.12
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xFA, 0xF3, 0xF0)
    sh.line.color.rgb = ORANGE
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    t = sh.text_frame
    t.word_wrap = True
    t.margin_left = t.margin_right = Emu(45000)
    t.margin_top = t.margin_bottom = Emu(28000)
    line(t, head, size=9, bold=True, colour=ORANGE, first=True)
    p = line(t, sub, size=7.5, colour=MUTED, space_before=1)
    for para in t.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    if i < len(steps) - 1:
        cx = x + box_w
        con = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, cx + Emu(28000),
                                  y + Emu(268000), Emu(94000), Emu(84000))
        con.fill.solid()
        con.fill.fore_color.rgb = ORANGE
        con.line.fill.background()
        con.shadow.inherit = False

# ---------------------------------------------------------------- slide 8
tf = body_of(S[7]).text_frame
tf.clear()
q(tf, 1, 'What alternatives exist today, and how does your solution compare?', first=True)
b(tf, 'Job boards (LinkedIn, Naukri, Internshala) — listings only; no resume help, no outcome memory.')
b(tf, 'Trackers and tailoring tools (Teal, Simplify) — closest comparison, but priced and positioned for the '
      'US market and still separate from interview preparation.')
b(tf, 'Resume builders (Canva, Zety) — formatting, not fit to a specific posting.')
b(tf, 'General chatbots — capable advice, but no memory of your resume, your applications or what already '
      'worked. The user re-explains themselves every session.')
a(tf, 'Navixa is the only one of these where a single context spans search, resume, tracker and interview prep.',
  space_before=4)
q(tf, 2, 'If a foundation-model provider shipped this feature tomorrow, why do you still win?')
a(tf, 'A model can generate advice; it cannot own the user’s outcome history. Our tracker records what '
      'happened after each application — which resume version reached an interview, which stalled, how long '
      'each stage took. That ties specific phrasing to real outcomes for specific roles and companies, and can '
      'only be accumulated by users completing full search cycles. It cannot be scraped or bought. Near-zero '
      'marginal infrastructure cost also lets us serve price-sensitive students at a price a funded competitor '
      'struggles to match.')

# ---------------------------------------------------------------- slide 9
tf = body_of(S[8]).text_frame
tf.clear()
q(tf, 1, 'What are your next major product and business milestones?', first=True)
a(tf, 'Product', size=10, colour=INK, space_before=3, bold=True)
b(tf, 'Feed the resume and tracker history into every AI prompt, so advice is personal rather than generic — '
      'the single largest quality gain available to us today.')
b(tf, 'Build an evaluation benchmark for career advice, then decide on fine-tuning against measured results '
      'rather than assumption.')
b(tf, 'Placement-cell dashboard: cohort view, outcome tracking, institutional reporting.')
a(tf, 'Business', size=10, colour=INK, space_before=6, bold=True)
b(tf, 'Pilot with two to three college placement cells; instrument outcomes from day one.')
b(tf, 'Convert one pilot into a paid institutional contract; introduce a low-priced individual tier.')
q(tf, 2, 'What is your long-term vision for the startup?')
a(tf, 'To become the system of record for a person’s career — every application, every resume version, '
      'every outcome in one place — so that career guidance becomes evidence-based and specific to the '
      'individual, instead of the generic advice most first-generation job seekers are left with.')

# ---------------------------------------------------------------- slide 10
tf = body_of(S[9]).text_frame
tf.clear()
q(tf, 1, 'Why is your team uniquely positioned to solve this problem?', first=True)
a(tf, 'We are the users. We are early-career job seekers in India running this exact search, which is why the '
      'product solves the workflow as it is actually lived rather than as imagined. The problems it fixes — '
      'rewriting the same resume for every posting, losing track of applications, generic interview prep — '
      'are ones we hit ourselves.')
q(tf, 2, 'What domain and technical expertise does the founding team bring?')
a(tf, 'Prajyot Kumar — Founder, Engineering', size=10, bold=True, colour=INK, space_before=4)
b(tf, 'Designed and shipped the entire product end-to-end: front end, serverless back end, database schema and '
      'row-level security, AI integration with fallbacks, CI-style test suite and deployment pipeline.')
b(tf, 'Ran a full security audit of the product, fixed two exploitable vulnerabilities and documented the '
      'residual risk honestly rather than claiming completeness.')
a(tf, '[Co-founder name] — [Role]', size=10, bold=True, colour=ORANGE, space_before=6)
b(tf, '[REPLACE BEFORE SUBMITTING: background, domain expertise, and what they own in the business.]')

# ---------------------------------------------------------------- slide 11
tf = body_of(S[10]).text_frame
tf.clear()
q(tf, 1, 'Why have you applied to the BITSoM Vertex programme?', first=True)
a(tf, 'We can build and ship — the product is live and tested. What we lack is commercial judgement and a route '
      'to market. Vertex offers the two things we cannot self-teach quickly: mentorship on turning a working '
      'product into a business, and institutional credibility with the colleges and placement cells that are '
      'our beachhead channel.')
q(tf, 2, 'Which challenge is currently limiting your startup’s growth?')
a(tf, 'Distribution and validation, not engineering. We have a working product and no reliable way to put it in '
      'front of users at scale. Placement cells are the right channel, but reaching them needs warm '
      'introductions and institutional trust that a student-built product does not have on its own. Until we '
      'are in front of real cohorts we also cannot gather the outcome data that makes the product defensible.')

# ---------------------------------------------------------------- slide 12
tf = body_of(S[11]).text_frame
tf.clear()
line(tf, 'Provide links to your:', size=10.5, bold=True, colour=INK, first=True)
rows = [('Deployed Project Link', 'https://navixa-woad.vercel.app', False),
        ('Project Demo Video (3–5 minutes)', '[ADD LINK BEFORE SUBMITTING]', True),
        ('Website', 'https://navixa-woad.vercel.app', False),
        ('GitHub Repository', 'https://github.com/prajyot2003/Navixa  (set to Public before submitting)', True),
        ('Product Documentation', 'README.md and SECURITY.md in the repository', False),
        ('Contact Details', 'Prajyot Kumar — prajyotkumar2003@gmail.com', False)]
for label, val, flag in rows:
    p = line(tf, '', size=10, indent=0.28, space_before=6, bullet=True)
    r1 = p.runs[0]
    r1.text = f'{label}:  '
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.name = 'Arial'
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(10)
    r2.font.name = 'Arial'
    r2.font.color.rgb = ORANGE if flag else BODY

prs.save(OUTPUT)
print(f'wrote {OUTPUT}')
